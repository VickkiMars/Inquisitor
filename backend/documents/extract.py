from typing import Optional
from backend.log_helper.report import log_message
import os

class Document:
    def __init__(self, file_path, num_questions):
        self.file_path=file_path
        self.num_questions = num_questions
        self.text = self.process_file()
        
    def process_file(self):
        _, ext = os.path.splitext(self.file_path)
        log_message(f"Processing file: {self.file_path}, {ext}")

        try:
            if ext == ".pdf":
                return self.extract_pdf()
            elif ext == ".docx":
                return self.extract_docx()
            elif ext == ".txt":
                return self.extract_txt()
            elif ext == ".epub":
                return self.extract_epub()
            elif ext == ".pptx":
                return self.extract_pptx()
            else:
                log_message(f"Unsupported file format: {ext}","error")
        except Exception as e:
            log_message(f"An unexpected error occurred: {e}", "error")
            return None
        
    def extract_pdf(self) -> str:
        log_message(f"Processing PDF: {self.file_path}")
        from PyPDF2 import PdfReader

        try:
            reader = PdfReader(self.file_path)
        except Exception as e:
            log_message(f"Could not process PDF: {e}", "error")

        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text.strip()

    def extract_docx(self) -> str:
        import docx
        log_message("Success importing 'docx': ")
        try:
            doc = docx.Document(self.file_path)
        except Exception as e:
            log_message(f"Could not process DOCX: {e}", "Error")
        return "\n".join(p.text for p in doc.paragraphs)

    def extract_txt(self) -> str:
        with open(self.file_path, "r", encoding="utf-8") as f:
            return f.read()
        
    def extract_epub(self) -> str:
        import ebooklib
        from ebooklib import epub

        book = epub.read_epub(self.file_path)
        text = ""
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                text += item.content.decode("utf-8")
        return text.strip()

    def extract_pptx(self) -> str:
        from pptx import Presentation
        prs = Presentation(self.file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        return text.strip()